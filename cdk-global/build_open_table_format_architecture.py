#!/usr/bin/env python3
"""Generate Databricks-branded Open Table Format architecture deck (python-pptx)."""

import ctypes.util
import os
import tempfile

_orig_find_library = ctypes.util.find_library


def _find_library_with_homebrew(name):
    result = _orig_find_library(name)
    if result is None and name in ("cairo", "cairo-2", "libcairo-2"):
        for path in (
            "/opt/homebrew/lib/libcairo.2.dylib",
            "/usr/local/lib/libcairo.2.dylib",
        ):
            if os.path.exists(path):
                return path
    return result


ctypes.util.find_library = _find_library_with_homebrew

try:
    import cairosvg

    HAS_CAIROSVG = True
except (ImportError, OSError):
    HAS_CAIROSVG = False

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BG_DARK = RGBColor(0x1B, 0x1F, 0x23)
BG_CARD = RGBColor(0x25, 0x2A, 0x30)
ACCENT_RED = RGBColor(0xFF, 0x38, 0x21)
ACCENT_GOLD = RGBColor(0xFF, 0xB7, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x00, 0xC8, 0x53)
BLUE = RGBColor(0x42, 0xA5, 0xF5)

ASSETS_DIR = (
    "/Users/zach.jacobson/Desktop/github_repos/customer_repos/presentation_automation/"
    "Databricks-selected-assets"
)

LOGO_MAP = {
    "databricks_main": os.path.join(
        ASSETS_DIR,
        "Databricks One Lockup Full Color",
        "databricks-one-lockup-full-color-white.svg",
    ),
    "delta_lake": os.path.join(ASSETS_DIR, "logo-color-delta-lake.svg"),
    "unity_catalog": os.path.join(
        ASSETS_DIR,
        "Unity Catalog Lockup Full Color",
        "unity-catalog-lockup-full-color-white.svg",
    ),
    "lakehouse": os.path.join(
        ASSETS_DIR,
        "Lakehouse Lockup Full Color",
        "lakehouse-lockup-full-color-white.svg",
    ),
    "delta_sharing": os.path.join(
        ASSETS_DIR,
        "Delta Sharing Lockup Full Color",
        "delta-sharing-lockup-full-color-white.svg",
    ),
    "mosaic_ai": os.path.join(
        ASSETS_DIR,
        "Mosaic AI Lockup Full Color",
        "mosaic-ai-lockup-full-color-white.svg",
    ),
}

_tmp_dir = tempfile.mkdtemp(prefix="db_logos_")


def _svg_to_png(svg_path, width_px=600):
    if not HAS_CAIROSVG or not os.path.exists(svg_path):
        return None
    png_path = os.path.join(_tmp_dir, os.path.basename(svg_path).replace(".svg", ".png"))
    if os.path.exists(png_path):
        return png_path
    try:
        cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=width_px)
        return png_path
    except Exception:
        return None


def _dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def _add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tx_box


def _add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _accent_bar(
    slide,
    left,
    top,
    width=Inches(0.08),
    height=Inches(0.6),
    color=ACCENT_RED,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _horizontal_rule(slide, left, top, width, color=ACCENT_RED, height=Inches(0.04)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_logo(slide, logo_key, left, top, width=None, height=None):
    svg_path = LOGO_MAP.get(logo_key)
    if not svg_path:
        return None
    png_path = _svg_to_png(svg_path)
    if not png_path:
        return None
    kwargs = {}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    if not kwargs:
        kwargs["height"] = Inches(0.6)
    return slide.shapes.add_picture(png_path, left, top, **kwargs)


def _add_watermark(slide):
    _add_logo(slide, "databricks_main", Inches(10.8), Inches(6.8), height=Inches(0.4))


def _section_header(slide, text, subtitle=None):
    _dark_bg(slide)
    _add_watermark(slide)
    _accent_bar(slide, Inches(0.6), Inches(1.6), Inches(0.08), Inches(0.8), ACCENT_RED)
    _add_text_box(
        slide,
        Inches(0.9),
        Inches(1.5),
        Inches(11),
        Inches(1),
        text,
        font_size=36,
        color=WHITE,
        bold=True,
    )
    if subtitle:
        _add_text_box(
            slide,
            Inches(0.9),
            Inches(2.4),
            Inches(11),
            Inches(0.7),
            subtitle,
            font_size=18,
            color=LIGHT_GRAY,
        )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Cost of dual-platform architecture
    s1 = prs.slides.add_slide(blank)
    _section_header(
        s1,
        "The Cost of Dual-Platform Data Architecture",
        "Snowflake + Databricks without open table formats",
    )
    _add_logo(s1, "lakehouse", Inches(10.5), Inches(1.45), height=Inches(0.5))
    _add_logo(s1, "delta_sharing", Inches(0.55), Inches(1.45), height=Inches(0.45))

    card_w = Inches(2.85)
    gap = Inches(0.25)
    start_x = Inches(0.5)
    y = Inches(3.15)
    h = Inches(2.35)
    cards = [
        (
            "Data Duplication",
            "Storing the same datasets in both platforms drives up storage costs.",
            ACCENT_RED,
        ),
        (
            "Compute Duplication",
            "Redundant queries and transformations run in Snowflake and Databricks.",
            BLUE,
        ),
        (
            "Networking Cost & Latency",
            "Moving data between platforms adds egress fees and delays.",
            ACCENT_GOLD,
        ),
        (
            "Dual Governance Overhead",
            "Policies, access controls, and lineage maintained separately in two systems.",
            GREEN,
        ),
    ]
    for i, (title, body, bar_color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        _add_card(s1, x, y, card_w, h)
        _accent_bar(s1, x + Inches(0.12), y + Inches(0.2), Inches(0.06), Inches(1.9), bar_color)
        _add_text_box(
            s1,
            x + Inches(0.28),
            y + Inches(0.22),
            card_w - Inches(0.42),
            Inches(0.45),
            title,
            font_size=15,
            color=WHITE,
            bold=True,
        )
        _add_text_box(
            s1,
            x + Inches(0.28),
            y + Inches(0.72),
            card_w - Inches(0.42),
            Inches(1.45),
            body,
            font_size=12,
            color=LIGHT_GRAY,
        )

    callout_y = Inches(5.75)
    _add_card(s1, Inches(0.5), callout_y, Inches(12.3), Inches(1.1), fill_color=RGBColor(0x2A, 0x32, 0x3A))
    _horizontal_rule(s1, Inches(0.65), callout_y + Inches(0.12), Inches(0.12), ACCENT_GOLD, Inches(0.85))
    _add_text_box(
        s1,
        Inches(0.85),
        callout_y + Inches(0.18),
        Inches(11.7),
        Inches(0.95),
        "These challenges compound as AI/ML use cases scale — making open table formats "
        "a prerequisite, not an optimization.",
        font_size=14,
        color=ACCENT_GOLD,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    # Slide 2 — Use cases blocked
    s2 = prs.slides.add_slide(blank)
    _section_header(
        s2,
        "Use Cases Blocked Without Open Table Formats",
        "Production readiness depends on a unified data layer",
    )
    _add_logo(s2, "mosaic_ai", Inches(10.5), Inches(1.45), height=Inches(0.5))
    _add_logo(s2, "unity_catalog", Inches(0.55), Inches(1.45), height=Inches(0.45))

    left_x = Inches(0.75)
    col_w = Inches(5.9)
    ty = Inches(3.05)
    use_cases = [
        ("1. AIVA 2.0", "All sub-agents need unified access across Snowflake and external systems."),
        ("2. MCP Gateway", "Requires a consistent data access layer across platforms."),
        ("3. WRAP", "Needs streamlined data pipelines without duplication."),
        ("4. Customer LTV", "Cross-platform analytics depend on unified data."),
        ("5. AI Platform", "Enterprise AI requires open, interoperable data as a foundation."),
        ("6. Parts Kit Agent", "Agent workflows need real-time, governed data access."),
    ]
    for idx, (title, desc) in enumerate(use_cases):
        col = idx % 2
        row = idx // 2
        x = left_x + col * (col_w + Inches(0.45))
        yy = ty + row * Inches(0.95)
        _add_text_box(s2, x, yy, col_w, Inches(0.32), title, font_size=15, color=WHITE, bold=True)
        _add_text_box(s2, x, yy + Inches(0.34), col_w, Inches(0.58), desc, font_size=12, color=LIGHT_GRAY)

    note_y = Inches(6.05)
    _add_card(s2, Inches(0.5), note_y, Inches(12.3), Inches(0.95), fill_color=BG_CARD)
    _accent_bar(s2, Inches(0.62), note_y + Inches(0.15), Inches(0.06), Inches(0.65), ACCENT_RED)
    _add_text_box(
        s2,
        Inches(0.78),
        note_y + Inches(0.18),
        Inches(11.85),
        Inches(0.72),
        "Each use case depends on a unified, open data layer — without it, teams face fragmented data, "
        "duplicated effort, and delayed time-to-production.",
        font_size=13,
        color=MED_GRAY,
    )

    # Slide 3 — Path forward
    s3 = prs.slides.add_slide(blank)
    _section_header(
        s3,
        "The Path Forward — Open Table Formats",
        "Iceberg and Delta as the interoperability standard",
    )
    _add_logo(s3, "delta_lake", Inches(10.2), Inches(1.4), height=Inches(0.55))
    _add_logo(s3, "unity_catalog", Inches(0.55), Inches(1.45), height=Inches(0.45))

    body_left = Inches(0.75)
    body_top = Inches(3.05)
    body_w = Inches(7.6)
    bullets = [
        "Prescribed architecture: adopt Iceberg (and/or Delta) as the open table format standard so "
        "Databricks can read directly from Snowflake-managed tables with minimal data movement.",
        "Start with the gold layer: implement incrementally — begin with curated/gold datasets to limit disruption.",
        "External client access: standardizing on Iceberg/Delta lets partners and clients read from Snowflake "
        "through open, portable interfaces.",
        "Customer references: major Databricks customers including UHG (UnitedHealth Group) and John Deere "
        "have adopted this interoperability pattern at scale.",
    ]
    for i, b in enumerate(bullets):
        _add_text_box(
            s3,
            body_left,
            body_top + i * Inches(0.78),
            body_w,
            Inches(0.75),
            f"• {b}",
            font_size=13,
            color=LIGHT_GRAY,
        )

    ref_x = Inches(8.55)
    ref_y = Inches(3.05)
    ref_w = Inches(4.0)
    _add_card(s3, ref_x, ref_y, ref_w, Inches(3.15))
    _accent_bar(s3, ref_x + Inches(0.12), ref_y + Inches(0.2), Inches(0.06), Inches(2.75), BLUE)
    _add_text_box(
        s3,
        ref_x + Inches(0.28),
        ref_y + Inches(0.22),
        ref_w - Inches(0.4),
        Inches(0.4),
        "Social proof",
        font_size=16,
        color=WHITE,
        bold=True,
    )
    _add_text_box(
        s3,
        ref_x + Inches(0.28),
        ref_y + Inches(0.68),
        ref_w - Inches(0.4),
        Inches(2.35),
        "UHG (UnitedHealth Group) and John Deere exemplify enterprises that have moved toward "
        "open table formats and governed lakehouse patterns to scale analytics and AI.",
        font_size=12,
        color=LIGHT_GRAY,
    )

    close_y = Inches(6.15)
    _horizontal_rule(s3, Inches(0.6), close_y, Inches(12.1), ACCENT_RED)
    _add_text_box(
        s3,
        Inches(0.75),
        close_y + Inches(0.12),
        Inches(11.8),
        Inches(0.85),
        "Open table formats aren't just a technical decision — they're the foundation for a scalable, "
        "governed, AI-ready data platform.",
        font_size=15,
        color=ACCENT_GOLD,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # ------------------------------------------------------------------
    # Slide 4 — Best Solution: Position UC as the Managing Catalog
    # ------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank)
    _dark_bg(s4)
    _add_watermark(s4)

    # "Best Solution" badge
    badge4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(1.6), Inches(0.4))
    badge4.fill.solid()
    badge4.fill.fore_color.rgb = GREEN
    badge4.line.fill.background()
    _add_text_box(s4, Inches(0.55), Inches(0.3), Inches(1.5), Inches(0.4),
                  "Best Solution", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_text_box(s4, Inches(0.5), Inches(0.85), Inches(12), Inches(0.6),
                  "Position UC as the Managing Catalog", font_size=30, color=WHITE, bold=True)
    _add_text_box(s4, Inches(0.5), Inches(1.45), Inches(12), Inches(0.4),
                  "Interop. Solution: Write directly to UC from SF", font_size=16, color=LIGHT_GRAY)

    # --- Left: Architecture flow (text-based) ---
    flow_x = Inches(0.5)
    flow_y = Inches(2.1)
    flow_w = Inches(5.8)
    _add_card(s4, flow_x, flow_y, flow_w, Inches(3.8))
    _add_text_box(s4, flow_x + Inches(0.2), flow_y + Inches(0.15), flow_w - Inches(0.4), Inches(0.35),
                  "Architecture Flow", font_size=14, color=WHITE, bold=True)
    _horizontal_rule(s4, flow_x + Inches(0.2), flow_y + Inches(0.5), Inches(2.0), ACCENT_RED, Inches(0.03))

    arch_lines = [
        "Unity Catalog (Metastore)  →  External Data Access + Grant External Use",
        "Snowflake Horizon (Account)  →  UC Catalog Integration + Linked Databases",
        "SF Warehouse  →  Write directly in UC from SF",
        "UC Managed Iceberg (UC Storage Root)  ←→  Read/Write from Databricks",
    ]
    for j, line in enumerate(arch_lines):
        _add_text_box(s4, flow_x + Inches(0.25), flow_y + Inches(0.65) + j * Inches(0.55),
                      flow_w - Inches(0.5), Inches(0.5), f"▸  {line}", font_size=11, color=LIGHT_GRAY)

    # Numbered steps below architecture
    steps4 = [
        "① Enable external data access to Unity Catalog + Grant external use schema to the service principal",
        "② Configure a catalog integration for Unity Catalog",
        "③ Create catalog-linked database",
        "④ Write directly to Databricks tables (AWS only at the moment)",
        "⑤ Read and Write to the same table from Databricks",
    ]
    step_y = flow_y + Inches(2.85)
    for j, step in enumerate(steps4):
        _add_text_box(s4, flow_x + Inches(0.15), step_y + j * Inches(0.18),
                      flow_w - Inches(0.3), Inches(0.2), step, font_size=8, color=MED_GRAY)

    # --- Right: Summary table as stacked cards ---
    tbl_x = Inches(6.55)
    tbl_y = Inches(2.1)
    tbl_w = Inches(6.3)

    # Summary card
    _add_card(s4, tbl_x, tbl_y, tbl_w, Inches(2.3))
    _accent_bar(s4, tbl_x + Inches(0.1), tbl_y + Inches(0.12), Inches(0.05), Inches(2.05), GREEN)
    _add_text_box(s4, tbl_x + Inches(0.25), tbl_y + Inches(0.08), Inches(1.0), Inches(0.3),
                  "Summary", font_size=12, color=WHITE, bold=True)
    summary_lines = [
        "★ BEST: When customer has Unity Catalog. Enables both Databricks and Snowflake "
        "to read and write on same Iceberg table managed by UC.",
        "▸ UC Managed Iceberg resides on external locations (S3/ADLS/GCS), not in Snowflake proprietary storage",
        "▸ Reads and Writes on same UC table from both Snowflake and Databricks",
        "▸ Customer gets Predictive Optimization benefits — auto-expires old snapshots, "
        "deletes unreferenced files, incrementally clusters via Liquid Clustering",
    ]
    for j, line in enumerate(summary_lines):
        _add_text_box(s4, tbl_x + Inches(0.25), tbl_y + Inches(0.38) + j * Inches(0.45),
                      tbl_w - Inches(0.4), Inches(0.42), line, font_size=9, color=LIGHT_GRAY)

    # Cost / Metadata / Governance row
    detail_y = tbl_y + Inches(2.45)
    detail_cards = [
        ("Cost", "✅ Writes supported from both engines\n✅ Cost depends on engine performing read/write", GREEN),
        ("Metadata", "✅ Updated automatically and incrementally during a read or refresh", BLUE),
        ("Governance", "Managed by Unity Catalog", ACCENT_GOLD),
    ]
    det_w = Inches(2.03)
    for j, (label, desc, bar_c) in enumerate(detail_cards):
        dx = tbl_x + j * (det_w + Inches(0.1))
        _add_card(s4, dx, detail_y, det_w, Inches(1.15))
        _accent_bar(s4, dx + Inches(0.08), detail_y + Inches(0.1), Inches(0.04), Inches(0.95), bar_c)
        _add_text_box(s4, dx + Inches(0.18), detail_y + Inches(0.08), det_w - Inches(0.25), Inches(0.25),
                      label, font_size=10, color=WHITE, bold=True)
        _add_text_box(s4, dx + Inches(0.18), detail_y + Inches(0.35), det_w - Inches(0.25), Inches(0.75),
                      desc, font_size=8, color=LIGHT_GRAY)

    # ------------------------------------------------------------------
    # Slide 5 — Best Solution: SF Managed Iceberg + UC Federation
    # ------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank)
    _dark_bg(s5)
    _add_watermark(s5)

    # Badges
    badge5 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(1.6), Inches(0.4))
    badge5.fill.solid()
    badge5.fill.fore_color.rgb = GREEN
    badge5.line.fill.background()
    _add_text_box(s5, Inches(0.55), Inches(0.3), Inches(1.5), Inches(0.4),
                  "Best Solution", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    badge5a = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.25), Inches(0.3), Inches(0.8), Inches(0.4))
    badge5a.fill.solid()
    badge5a.fill.fore_color.rgb = BLUE
    badge5a.line.fill.background()
    _add_text_box(s5, Inches(2.3), Inches(0.3), Inches(0.7), Inches(0.4),
                  "AWS", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_text_box(s5, Inches(0.5), Inches(0.85), Inches(12), Inches(0.6),
                  "Move to Snowflake Managed Iceberg from Native Tables", font_size=28, color=WHITE, bold=True)
    _add_text_box(s5, Inches(0.5), Inches(1.45), Inches(12), Inches(0.4),
                  "Interop. Solution: Unity Catalog Federation to SF Horizon", font_size=16, color=LIGHT_GRAY)

    # --- Left: Architecture flow ---
    _add_card(s5, flow_x, flow_y, flow_w, Inches(3.8))
    _add_text_box(s5, flow_x + Inches(0.2), flow_y + Inches(0.15), flow_w - Inches(0.4), Inches(0.35),
                  "Architecture Flow", font_size=14, color=WHITE, bold=True)
    _horizontal_rule(s5, flow_x + Inches(0.2), flow_y + Inches(0.5), Inches(2.0), ACCENT_RED, Inches(0.03))

    arch_lines5 = [
        "Unity Catalog (Metastore)  →  Connection to Snowflake",
        "External Location  →  Points to Snowflake storage (S3/ADLS/GCS)",
        "Foreign Catalog  →  Created in UC using the connection + external location",
        "Databricks Workspace  →  Queries SF tables via UC Federation",
        "Object Storage (Managed Iceberg)  ←→  SF Warehouse writes, DBX reads from storage",
    ]
    for j, line in enumerate(arch_lines5):
        _add_text_box(s5, flow_x + Inches(0.25), flow_y + Inches(0.65) + j * Inches(0.48),
                      flow_w - Inches(0.5), Inches(0.45), f"▸  {line}", font_size=11, color=LIGHT_GRAY)

    # Numbered steps
    steps5 = [
        "① Create a connection to Snowflake on Unity Catalog",
        "② Create an External Location in UC pointing to Snowflake storage (create Storage Root location)",
        "③ Create a foreign catalog in UC — add External Location as authorized path + specify external storage root",
        "④ Query your Snowflake table in Databricks — Iceberg tables read directly from storage; "
        "native SF tables fall back to JDBC",
    ]
    step_y5 = flow_y + Inches(3.05)
    for j, step in enumerate(steps5):
        _add_text_box(s5, flow_x + Inches(0.15), step_y5 + j * Inches(0.19),
                      flow_w - Inches(0.3), Inches(0.22), step, font_size=8, color=MED_GRAY)

    # --- Right: Summary table ---
    _add_card(s5, tbl_x, tbl_y, tbl_w, Inches(2.55))
    _accent_bar(s5, tbl_x + Inches(0.1), tbl_y + Inches(0.12), Inches(0.05), Inches(2.3), GREEN)
    _add_text_box(s5, tbl_x + Inches(0.25), tbl_y + Inches(0.08), Inches(1.0), Inches(0.3),
                  "Summary", font_size=12, color=WHITE, bold=True)
    summary5 = [
        "★ Move native SF tables to Snowflake managed Iceberg tables",
        "▸ SF Managed Iceberg resides on external volumes (S3/ADLS/GCS), not Snowflake proprietary storage — "
        "UC can access it and avoid double compute",
        "▸ Read-only in Databricks",
        "⚠ SF Iceberg table read will fall back to JDBC if not present under authorized paths "
        "(costly due to double compute)",
        "⚠ When creating a UC catalog, specify an external storage root — default root has known issues",
    ]
    for j, line in enumerate(summary5):
        _add_text_box(s5, tbl_x + Inches(0.25), tbl_y + Inches(0.38) + j * Inches(0.40),
                      tbl_w - Inches(0.4), Inches(0.38), line, font_size=9, color=LIGHT_GRAY)

    # Cost / Metadata / Governance row
    detail_y5 = tbl_y + Inches(2.7)
    detail_cards5 = [
        ("Cost",
         "✅ Single Compute on DBX side\n"
         "✅ Very small SF WH for Iceberg metadata\n"
         "✅ No data replication — efficient storage", GREEN),
        ("Metadata",
         "✅ Updated automatically and incrementally during a read or refresh", BLUE),
        ("Governance",
         "Managed by Snowflake Horizon; access controls also need to be applied in UC", ACCENT_GOLD),
    ]
    for j, (label, desc, bar_c) in enumerate(detail_cards5):
        dx = tbl_x + j * (det_w + Inches(0.1))
        _add_card(s5, dx, detail_y5, det_w, Inches(1.15))
        _accent_bar(s5, dx + Inches(0.08), detail_y5 + Inches(0.1), Inches(0.04), Inches(0.95), bar_c)
        _add_text_box(s5, dx + Inches(0.18), detail_y5 + Inches(0.08), det_w - Inches(0.25), Inches(0.25),
                      label, font_size=10, color=WHITE, bold=True)
        _add_text_box(s5, dx + Inches(0.18), detail_y5 + Inches(0.35), det_w - Inches(0.25), Inches(0.75),
                      desc, font_size=8, color=LIGHT_GRAY)

    out = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "Open_Table_Format_Architecture.pptx",
    )
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
