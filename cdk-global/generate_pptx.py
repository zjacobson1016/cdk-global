#!/usr/bin/env python3
"""Generate a PowerPoint presentation from the CDK Global README content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Databricks brand-ish palette ──
BG_DARK     = RGBColor(0x1B, 0x1F, 0x23)   # dark charcoal
BG_CARD     = RGBColor(0x25, 0x2A, 0x30)   # card bg
ACCENT_RED  = RGBColor(0xFF, 0x38, 0x21)   # Databricks red
ACCENT_GOLD = RGBColor(0xFF, 0xB7, 0x2B)   # gold highlight
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xBB, 0xBB, 0xBB)
MED_GRAY    = RGBColor(0x88, 0x88, 0x88)
GREEN       = RGBColor(0x00, 0xC8, 0x53)
BLUE        = RGBColor(0x42, 0xA5, 0xF5)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Helper functions ──

def _dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, bullet_color=ACCENT_RED):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def _add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _section_header(slide, text, subtitle=None):
    _dark_bg(slide)
    _accent_bar(slide, Inches(0.6), Inches(1.6), Inches(0.08), Inches(0.8), ACCENT_RED)
    _add_text_box(slide, Inches(0.9), Inches(1.5), Inches(11), Inches(1),
                  text, font_size=36, color=WHITE, bold=True)
    if subtitle:
        _add_text_box(slide, Inches(0.9), Inches(2.4), Inches(11), Inches(0.7),
                      subtitle, font_size=18, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_dark_bg(slide)

# Red accent line
_accent_bar(slide, Inches(0.6), Inches(2.2), Inches(1.5), Inches(0.06), ACCENT_RED)

_add_text_box(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(1.2),
              "CDK Global", font_size=52, color=WHITE, bold=True)
_add_text_box(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(1),
              "Databricks Solution Accelerator", font_size=32, color=ACCENT_GOLD, bold=True)
_add_text_box(slide, Inches(0.6), Inches(4.5), Inches(10), Inches(1),
              "End-to-end auto lending: MLOps  |  Document AI  |  Deterministic Rules  |  Multi-Agent Supervisor",
              font_size=18, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 – Solution Overview
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Solution Overview",
                "Three specialized agents orchestrated by a Multi-Agent Supervisor")

# ── MAS card (top center) ──
_add_card(slide, Inches(3.5), Inches(3.2), Inches(6.3), Inches(0.9), ACCENT_RED)
_add_text_box(slide, Inches(3.7), Inches(3.3), Inches(5.9), Inches(0.7),
              "CDK Lending Supervisor (MAS)\nUnified conversational interface for F&I managers",
              font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ── Three agent cards ──
card_y = Inches(4.5)
card_h = Inches(2.2)
card_w = Inches(3.6)
gap = Inches(0.35)
start_x = Inches(1.2)

# Agent 1 – Lender Approval
x1 = start_x
_add_card(slide, x1, card_y, card_w, card_h)
_accent_bar(slide, x1 + Inches(0.15), card_y + Inches(0.15), Inches(0.06), Inches(0.5), ACCENT_RED)
_add_text_box(slide, x1 + Inches(0.35), card_y + Inches(0.1), card_w - Inches(0.5), Inches(0.5),
              "Lender Approval Agent", font_size=16, color=ACCENT_GOLD, bold=True)
_add_text_box(slide, x1 + Inches(0.35), card_y + Inches(0.55), card_w - Inches(0.5), Inches(1.5),
              "Model Serving Endpoint\n\nML model + deterministic rules\nIncome verification\nID expiration check\nStructured decision reasoning",
              font_size=12, color=LIGHT_GRAY)

# Agent 2 – Lending Analytics
x2 = start_x + card_w + gap
_add_card(slide, x2, card_y, card_w, card_h)
_accent_bar(slide, x2 + Inches(0.15), card_y + Inches(0.15), Inches(0.06), Inches(0.5), BLUE)
_add_text_box(slide, x2 + Inches(0.35), card_y + Inches(0.1), card_w - Inches(0.5), Inches(0.5),
              "Lending Analytics Genie", font_size=16, color=ACCENT_GOLD, bold=True)
_add_text_box(slide, x2 + Inches(0.35), card_y + Inches(0.55), card_w - Inches(0.5), Inches(1.5),
              "Genie Space (SQL)\n\nSelf-service analytics\nGold layer + inference tables\nApproval rates & trends\nNatural language queries",
              font_size=12, color=LIGHT_GRAY)

# Agent 3 – Lender Shopping
x3 = start_x + 2 * (card_w + gap)
_add_card(slide, x3, card_y, card_w, card_h)
_accent_bar(slide, x3 + Inches(0.15), card_y + Inches(0.15), Inches(0.06), Inches(0.5), GREEN)
_add_text_box(slide, x3 + Inches(0.35), card_y + Inches(0.1), card_w - Inches(0.5), Inches(0.5),
              "Lender Shopping Agent", font_size=16, color=ACCENT_GOLD, bold=True)
_add_text_box(slide, x3 + Inches(0.35), card_y + Inches(0.55), card_w - Inches(0.5), Inches(1.5),
              "Custom Agent Endpoint\n\nUC Function: shop_lenders()\nMulti-lender rate comparison\n8 lenders, 20 programs\nEstimated payments & likelihood",
              font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 – Data Pipeline Architecture
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Data Pipeline Architecture",
                "Three data sources joined by application_id through a medallion pipeline")

row_y = Inches(3.3)
box_w = Inches(3.4)
box_h = Inches(1.1)
gap_x = Inches(0.4)
start = Inches(1.5)

sources = [
    ("Structured Applications", "Parquet files\nAuto Loader ingestion", ACCENT_RED),
    ("Pay Stubs (PDF)", "ai_parse_document + ai_query\nVerified income extraction", BLUE),
    ("Photo IDs (JPEG)", "ai_parse_document + ai_query\nIdentity + expiration data", GREEN),
]
for i, (title, desc, color) in enumerate(sources):
    x = start + i * (box_w + gap_x)
    _add_card(slide, x, row_y, box_w, box_h)
    _accent_bar(slide, x + Inches(0.1), row_y + Inches(0.1), Inches(0.06), Inches(0.4), color)
    _add_text_box(slide, x + Inches(0.3), row_y + Inches(0.05), box_w - Inches(0.4), Inches(0.4),
                  title, font_size=15, color=ACCENT_GOLD, bold=True)
    _add_text_box(slide, x + Inches(0.3), row_y + Inches(0.45), box_w - Inches(0.4), Inches(0.6),
                  desc, font_size=12, color=LIGHT_GRAY)

# Arrow → medallion
_add_text_box(slide, Inches(5.5), Inches(4.5), Inches(2.5), Inches(0.5),
              "▼  JOIN on application_id  ▼", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Medallion row
medal_y = Inches(5.1)
medal_w = Inches(2.8)
medals = [
    ("Bronze", "3 streaming tables\n(apps, pay stubs, photo IDs)", ACCENT_RED),
    ("Silver", "LEFT JOIN all sources\nEnriched applications", ACCENT_GOLD),
    ("Gold", "ML features + verification\nsignals + train/test split", GREEN),
]
for i, (title, desc, color) in enumerate(medals):
    x = Inches(1.8) + i * (medal_w + Inches(0.5))
    _add_card(slide, x, medal_y, medal_w, Inches(1.2))
    _add_text_box(slide, x + Inches(0.15), medal_y + Inches(0.05), medal_w - Inches(0.3), Inches(0.4),
                  title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, x + Inches(0.15), medal_y + Inches(0.45), medal_w - Inches(0.3), Inches(0.7),
                  desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 – ML Pipeline + Business Rules
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "ML Pipeline + Deterministic Business Rules",
                "Optuna HPO with custom pyfunc wrapper that applies hard rules at prediction time")

# Left side – ML pipeline
_add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8))
_add_text_box(slide, Inches(0.9), Inches(3.3), Inches(5.2), Inches(0.5),
              "ML Pipeline", font_size=20, color=ACCENT_GOLD, bold=True)
steps = [
    "Feature Engineering: Feature table + 3 on-demand UC functions",
    "Optuna HPO: LogisticRegression, RandomForest, LightGBM",
    "Custom pyfunc wrapper (LenderApprovalWithRules)",
    "Champion / Challenger promotion flow",
    "Batch inference via fe.score_batch",
    "Real-time serving via Model Serving endpoint",
    "Lakehouse Monitoring + auto-retrain on drift",
]
_add_bullet_list(slide, Inches(0.9), Inches(3.85), Inches(5.2), Inches(2.8),
                 [f"  {s}" for s in steps], font_size=13, color=LIGHT_GRAY)

# Right side – Business rules
_add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8))
_add_text_box(slide, Inches(7.1), Inches(3.3), Inches(5.2), Inches(0.5),
              "Deterministic Rules", font_size=20, color=ACCENT_RED, bold=True)

_add_text_box(slide, Inches(7.1), Inches(3.9), Inches(5.2), Inches(0.4),
              "Rule 1: Income Validation", font_size=15, color=WHITE, bold=True)
_add_text_box(slide, Inches(7.1), Inches(4.25), Inches(5.2), Inches(0.8),
              "Verified income (pay stub x 26) must be\n70%-150% of self-reported income\nFAIL = automatic denial",
              font_size=13, color=LIGHT_GRAY)

_add_text_box(slide, Inches(7.1), Inches(5.1), Inches(5.2), Inches(0.4),
              "Rule 2: ID Expiration Check", font_size=15, color=WHITE, bold=True)
_add_text_box(slide, Inches(7.1), Inches(5.45), Inches(5.2), Inches(0.8),
              "Photo ID must not be expired\nFAIL = automatic denial",
              font_size=13, color=LIGHT_GRAY)

_add_text_box(slide, Inches(7.1), Inches(6.15), Inches(5.2), Inches(0.6),
              "Rules fire at both batch and real-time serving.\nEvery response includes decision_reason.",
              font_size=12, color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 – Prediction Response Format
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Structured Prediction Output",
                "Every prediction (batch & real-time) returns full decision reasoning")

_add_card(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(3.5))

json_text = (
    '{\n'
    '    "prediction": 0,\n'
    '    "ml_prediction": 1,\n'
    '    "ml_probability": 0.8723,\n'
    '    "income_check": "FAIL",\n'
    '    "id_check": "PASS",\n'
    '    "decision_reason": "DENIED by rules: Income mismatch\n'
    '                        (pay stub vs application)"\n'
    '}'
)
_add_text_box(slide, Inches(2.0), Inches(3.5), Inches(5), Inches(3),
              json_text, font_size=16, color=GREEN, font_name="Courier New")

annotations = [
    ("prediction", "Final decision after rules (0 = denied)"),
    ("ml_prediction", "Raw ML model output (1 = would approve)"),
    ("ml_probability", "Model confidence score"),
    ("income_check", "PASS / FAIL / MISSING"),
    ("id_check", "PASS / FAIL / MISSING"),
    ("decision_reason", "Human-readable explanation"),
]
y = Inches(3.5)
for field, desc in annotations:
    _add_text_box(slide, Inches(7.8), y, Inches(4.5), Inches(0.4),
                  f"{field}  {desc}", font_size=12, color=LIGHT_GRAY)
    y += Inches(0.4)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 – Lender Shopping Agent
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Lender Shopping Agent",
                "UC Function tool for multi-lender rate comparison")

# Function signature card
_add_card(slide, Inches(0.6), Inches(3.3), Inches(6.0), Inches(3.5))
_add_text_box(slide, Inches(0.9), Inches(3.4), Inches(5.5), Inches(0.4),
              "shop_lenders() UC Function", font_size=16, color=ACCENT_GOLD, bold=True)

sig_text = (
    "shop_lenders(\n"
    "    credit_score     INT,\n"
    "    annual_income    DOUBLE,\n"
    "    loan_amount      DOUBLE,\n"
    "    loan_term_months INT,\n"
    "    vehicle_year     INT\n"
    ")\n"
    "RETURNS TABLE(\n"
    "    lender_name, program_name, apr,\n"
    "    estimated_monthly_payment,\n"
    "    max_ltv, approval_likelihood\n"
    ")"
)
_add_text_box(slide, Inches(0.9), Inches(3.85), Inches(5.5), Inches(2.8),
              sig_text, font_size=13, color=GREEN, font_name="Courier New")

# Lenders card
_add_card(slide, Inches(6.9), Inches(3.3), Inches(5.8), Inches(3.5))
_add_text_box(slide, Inches(7.2), Inches(3.4), Inches(5.2), Inches(0.4),
              "8 Lenders  |  20 Programs", font_size=16, color=ACCENT_GOLD, bold=True)

lenders = [
    "TD Auto Finance          2.99% - 7.49%",
    "Chase Auto               3.49% - 9.49%",
    "Navy Federal CU          3.79% - 8.99%",
    "Ally Financial            3.99% - 21.99%",
    "Capital One Auto         4.49% - 14.99%",
    "Wells Fargo Dealer       4.29% - 17.99%",
    "AmeriCredit (GM Fin)     4.99% - 17.49%",
    "Westlake Financial      14.49% - 24.99%",
]
_add_bullet_list(slide, Inches(7.2), Inches(3.85), Inches(5.2), Inches(2.8),
                 lenders, font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 – Lending Analytics Genie Space
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Lending Analytics Genie Space",
                "Self-service SQL analytics via natural language")

# Tables card
_add_card(slide, Inches(0.6), Inches(3.3), Inches(5.8), Inches(3.5))
_add_text_box(slide, Inches(0.9), Inches(3.4), Inches(5.2), Inches(0.4),
              "Connected Tables", font_size=16, color=ACCENT_GOLD, bold=True)
tables = [
    "gold_lender_features  --  ML-ready features",
    "lender_approval_inference_table  --  predictions + labels",
    "lender_approval_offline_inference  --  decisions + reasoning",
    "lender_programs  --  lender rate reference data",
    "lender_approval_feature_table  --  feature store",
]
_add_bullet_list(slide, Inches(0.9), Inches(3.85), Inches(5.2), Inches(2.8),
                 tables, font_size=13, color=LIGHT_GRAY)

# Example questions card
_add_card(slide, Inches(6.9), Inches(3.3), Inches(5.8), Inches(3.5))
_add_text_box(slide, Inches(7.2), Inches(3.4), Inches(5.2), Inches(0.4),
              "Example Questions", font_size=16, color=ACCENT_GOLD, bold=True)
questions = [
    '"What is the overall approval rate this month?"',
    '"How many apps denied by rules vs ML model?"',
    '"Top 5 lenders by lowest APR for prime borrowers"',
    '"% of apps with income verification mismatches"',
    '"Compare approval rates by credit score tier"',
    '"Show denied apps where ML would have approved"',
]
_add_bullet_list(slide, Inches(7.2), Inches(3.85), Inches(5.2), Inches(2.8),
                 questions, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 – MAS Routing
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Multi-Agent Supervisor Routing",
                "Automatic question routing to the right specialist")

# Three routing rules
rules = [
    ("Approval / denial decisions", "lender_approval", ACCENT_RED,
     '"Should I approve this? Credit 680, income $52K, loan $28K"'),
    ("Data, metrics, trends, reports", "lending_analytics (Genie)", BLUE,
     '"What was our approval rate last week?"'),
    ("Rate comparisons, programs, payments", "lender_shopping", GREEN,
     '"Best rates for 740 credit, $72K income, $38K loan on 2025 Honda"'),
]

y = Inches(3.3)
for question_type, agent_name, color, example in rules:
    _add_card(slide, Inches(0.8), y, Inches(11.7), Inches(1.1))
    _accent_bar(slide, Inches(1.0), y + Inches(0.15), Inches(0.06), Inches(0.8), color)
    _add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(4.5), Inches(0.4),
                  question_type, font_size=16, color=WHITE, bold=True)
    _add_text_box(slide, Inches(1.3), y + Inches(0.45), Inches(4.5), Inches(0.5),
                  f"Routes to: {agent_name}", font_size=13, color=color, bold=True)
    _add_text_box(slide, Inches(6.5), y + Inches(0.2), Inches(5.7), Inches(0.7),
                  example, font_size=13, color=MED_GRAY)
    y += Inches(1.3)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 – Business Value
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Business Value for CDK Global")

metrics = [
    ("~166,000 hrs/year", "Underwriter hours saved", "At 500K apps, document review\ngoes from ~20 min to seconds", ACCENT_RED),
    ("$30-50M/year", "Fraud loss prevention", "Catch 1% of income-misrepresented\nloans on $15B portfolio", ACCENT_GOLD),
    ("100%", "Audit trail coverage", "Full lineage from raw document\nto final decision", BLUE),
    ("2-3% F1 lift", "Model accuracy improvement", "Continuous retraining via\ndrift detection", GREEN),
    ("Minutes to seconds", "Dealer decision latency", "Real-time serving +\nrate shopping", WHITE),
]

card_w = Inches(2.2)
card_h = Inches(2.8)
gap = Inches(0.3)
start_x = Inches(0.6)
y = Inches(3.3)

for i, (value, label, desc, color) in enumerate(metrics):
    x = start_x + i * (card_w + gap)
    _add_card(slide, x, y, card_w, card_h)
    _add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), Inches(0.7),
                  value, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, x + Inches(0.15), y + Inches(0.85), card_w - Inches(0.3), Inches(0.5),
                  label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, x + Inches(0.15), y + Inches(1.45), card_w - Inches(0.3), Inches(1.2),
                  desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 – Day 1 Schedule
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Day 1: Welcome & Platform Overview",
                "9:00 AM - 1:30 PM  |  Introductions, solution context, and Databricks platform deep dive")

col_time_x  = Inches(0.6)
col_sess_x  = Inches(2.6)
col_detail_x = Inches(7.0)
header_y = Inches(3.1)

_add_text_box(slide, col_time_x, header_y, Inches(1.8), Inches(0.4),
              "Time", font_size=13, color=MED_GRAY, bold=True)
_add_text_box(slide, col_sess_x, header_y, Inches(4.2), Inches(0.4),
              "Session", font_size=13, color=MED_GRAY, bold=True)
_add_text_box(slide, col_detail_x, header_y, Inches(5.5), Inches(0.4),
              "Key Topics", font_size=13, color=MED_GRAY, bold=True)

schedule_day1 = [
    ("9:00 - 9:30",   "Welcome & Solution Overview",   "All",      WHITE,
     "Introductions  |  Workshop Goals\nCDK Solution Architecture Walkthrough\nEnd-to-end demo preview"),
    ("9:30 - 10:45",  "Databricks Platform Overview",  "Platform", BLUE,
     "Lakehouse Architecture  |  Unity Catalog\nDelta Lake  |  Serverless Compute\nWorkspace, Clusters, SQL Warehouses"),
    ("10:45 - 11:00", "Break",                         "",         MED_GRAY, None),
    ("11:00 - 12:15", "Platform Deep Dive",            "Platform", BLUE,
     "Data Engineering: Auto Loader, SDP Pipelines\nSQL Analytics: Warehouses, Dashboards, Genie\nAI & GenAI: Foundation Models, Vector Search\nGovernance: Lineage, Row/Column Security"),
    ("12:15 - 12:45", "Lunch",                         "",         MED_GRAY, None),
    ("12:45 - 1:30",  "Platform Hands-On & Q&A",       "Platform", BLUE,
     "Live platform walkthrough\nWorkspace navigation & notebooks\nUnity Catalog exploration"),
]

y = Inches(3.5)
for time_str, session, track, color, details in schedule_day1:
    is_break = session in ("Break", "Lunch")
    row_h = Inches(0.35) if is_break else Inches(0.95) if details else Inches(0.5)

    if not is_break and details:
        _add_card(slide, Inches(0.5), y - Inches(0.05), Inches(12.3), row_h + Inches(0.1))
    if details:
        _accent_bar(slide, Inches(0.65), y + Inches(0.05), Inches(0.05), row_h - Inches(0.1), color)

    _add_text_box(slide, col_time_x + Inches(0.2), y, Inches(1.6), Inches(0.35),
                  time_str, font_size=12, color=LIGHT_GRAY if is_break else WHITE, bold=not is_break)
    _add_text_box(slide, col_sess_x, y, Inches(4.2), Inches(0.35),
                  session, font_size=13 if not is_break else 11,
                  color=MED_GRAY if is_break else color, bold=not is_break)

    if track and not is_break:
        _add_text_box(slide, col_sess_x, y + Inches(0.3), Inches(4.2), Inches(0.3),
                      f"Track: {track}", font_size=11, color=MED_GRAY)
    if details:
        _add_text_box(slide, col_detail_x, y + Inches(0.02), Inches(5.5), row_h - Inches(0.05),
                      details, font_size=11, color=LIGHT_GRAY)

    y += row_h + Inches(0.08)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10b – Day 2 Schedule
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Day 2: Hands-On Breakout Sessions",
                "9:00 AM - 1:00 PM  |  AI Engineering, Data Engineering, ML & MLOps/LLMOps")

col_time_x  = Inches(0.6)
col_sess_x  = Inches(2.6)
col_detail_x = Inches(7.0)
header_y = Inches(3.1)

_add_text_box(slide, col_time_x, header_y, Inches(1.8), Inches(0.4),
              "Time", font_size=13, color=MED_GRAY, bold=True)
_add_text_box(slide, col_sess_x, header_y, Inches(4.2), Inches(0.4),
              "Session", font_size=13, color=MED_GRAY, bold=True)
_add_text_box(slide, col_detail_x, header_y, Inches(5.5), Inches(0.4),
              "Key Topics", font_size=13, color=MED_GRAY, bold=True)

schedule_day2 = [
    ("9:00 - 10:15",  "AI Engineering",        "AI",             GREEN,
     "Agent Bricks  |  Genie Spaces\nUC Functions  |  Model Serving\nMulti-Agent Supervisor orchestration"),
    ("10:15 - 10:30", "Break",                  "",               MED_GRAY, None),
    ("10:30 - 11:45", "Data Engineering",       "Data Pipelines", ACCENT_RED,
     "Auto Loader  |  Document AI (ai_parse_document)\nSDP Pipeline: Bronze / Silver / Gold\n3-way join on application_id"),
    ("11:45 - 12:00", "Break",                  "",               MED_GRAY, None),
    ("12:00 - 1:00",  "ML & MLOps/LLMOps",     "ML / MLOps",     ACCENT_GOLD,
     "Feature Engineering + UC Functions\nOptuna HPO + pyfunc wrapper with rules\nChampion/Challenger  |  Serving  |  Monitoring"),
]

y = Inches(3.5)
for time_str, session, track, color, details in schedule_day2:
    is_break = session == "Break"
    row_h = Inches(0.35) if is_break else Inches(0.95) if details else Inches(0.5)

    if not is_break and details:
        _add_card(slide, Inches(0.5), y - Inches(0.05), Inches(12.3), row_h + Inches(0.1))
    if details:
        _accent_bar(slide, Inches(0.65), y + Inches(0.05), Inches(0.05), row_h - Inches(0.1), color)

    _add_text_box(slide, col_time_x + Inches(0.2), y, Inches(1.6), Inches(0.35),
                  time_str, font_size=12, color=LIGHT_GRAY if is_break else WHITE, bold=not is_break)
    _add_text_box(slide, col_sess_x, y, Inches(4.2), Inches(0.35),
                  session, font_size=13 if not is_break else 11,
                  color=MED_GRAY if is_break else color, bold=not is_break)

    if track and not is_break:
        _add_text_box(slide, col_sess_x, y + Inches(0.3), Inches(4.2), Inches(0.3),
                      f"Track: {track}", font_size=11, color=MED_GRAY)
    if details:
        _add_text_box(slide, col_detail_x, y + Inches(0.02), Inches(5.5), row_h - Inches(0.05),
                      details, font_size=11, color=LIGHT_GRAY)

    y += row_h + Inches(0.08)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 – Databricks Platform Overview
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Databricks Platform Overview",
                "A unified analytics platform for data, AI, and governance  |  45 min")

# Left column – Lakehouse Architecture
_add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(0.8), Inches(3.35), Inches(0.05), Inches(0.5), BLUE)
_add_text_box(slide, Inches(1.05), Inches(3.3), Inches(5.1), Inches(0.5),
              "Lakehouse Architecture", font_size=18, color=BLUE, bold=True)
platform_items = [
    "Data Lakehouse: combines data lake + data warehouse",
    "Delta Lake: ACID transactions, time travel, Z-order/Liquid Clustering",
    "Unity Catalog: unified governance for tables, volumes, models, functions",
    "Serverless compute: instant-on SQL warehouses & job clusters",
    "Workflows: orchestrate notebooks, pipelines, and jobs",
    "Spark Declarative Pipelines: managed ETL (bronze/silver/gold)",
    "Photon: vectorized query engine for fast SQL & Spark",
]
_add_bullet_list(slide, Inches(1.05), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in platform_items], font_size=13, color=LIGHT_GRAY)

# Right column – Key capabilities
_add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(7.0), Inches(3.35), Inches(0.05), Inches(0.5), BLUE)
_add_text_box(slide, Inches(7.25), Inches(3.3), Inches(5.1), Inches(0.5),
              "Key Capabilities", font_size=18, color=BLUE, bold=True)
capabilities = [
    "Data Engineering: Auto Loader, SDP, Structured Streaming",
    "Data Science & ML: MLflow, Feature Store, Model Serving",
    "SQL Analytics: SQL Warehouses, AI/BI Dashboards, Genie",
    "AI & GenAI: Foundation Models, Agent Bricks, Vector Search",
    "Governance: Unity Catalog, Lineage, Row/Column Security",
    "Workspace: Notebooks, Repos, Git integration, CLI & SDK",
    "Marketplace: shared data products & solution accelerators",
]
_add_bullet_list(slide, Inches(7.25), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in capabilities], font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 – Deep Dive: Data Engineering
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Data Engineering",
                "Ingesting structured and unstructured data through a medallion architecture  |  75 min")

# Left column – What we build
_add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(0.8), Inches(3.35), Inches(0.05), Inches(0.5), ACCENT_RED)
_add_text_box(slide, Inches(1.05), Inches(3.3), Inches(5.1), Inches(0.5),
              "What We Build", font_size=18, color=ACCENT_RED, bold=True)
steps_de = [
    "Generate synthetic applications (parquet), pay stubs (PDF), photo IDs (JPEG)",
    "Store raw files in Unity Catalog Volumes",
    "Auto Loader ingestion (parquet + binaryFile)",
    "AI extraction: ai_parse_document + ai_query",
    "Bronze / Silver / Gold SDP pipeline",
    "3-way LEFT JOIN on application_id",
    "Derived ML features with verification signals",
]
_add_bullet_list(slide, Inches(1.05), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in steps_de], font_size=13, color=LIGHT_GRAY)

# Right column – Notebooks
_add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(7.0), Inches(3.35), Inches(0.05), Inches(0.5), ACCENT_RED)
_add_text_box(slide, Inches(7.25), Inches(3.3), Inches(5.1), Inches(0.5),
              "Notebooks & Files", font_size=18, color=ACCENT_RED, bold=True)
notebooks_de = [
    "01_generate_lender_data.py",
    "01_generate_lender_data_pdf.py",
    "bronze_applications.py    (Auto Loader parquet)",
    "bronze_pay_stubs.py       (AI extraction PDF)",
    "bronze_photo_ids.py       (AI extraction JPEG)",
    "silver_applications.py    (3-way join)",
    "gold_features.py          (ML features)",
]
_add_bullet_list(slide, Inches(7.25), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in notebooks_de], font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 – Deep Dive: ML & MLOps/LLMOps
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "ML & MLOps/LLMOps",
                "Feature engineering, model training with business rules, and production deployment  |  60 min")

_add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(0.8), Inches(3.35), Inches(0.05), Inches(0.5), ACCENT_GOLD)
_add_text_box(slide, Inches(1.05), Inches(3.3), Inches(5.1), Inches(0.5),
              "What We Build", font_size=18, color=ACCENT_GOLD, bold=True)
steps_ml = [
    "Feature table + 3 on-demand UC functions",
    "Optuna HPO (LogReg, RF, LightGBM)",
    "Custom pyfunc wrapper with deterministic rules",
    "Structured output: prediction + reasoning",
    "Champion / Challenger promotion flow",
    "Batch inference via fe.score_batch",
    "Real-time Model Serving endpoint",
    "Lakehouse Monitoring + auto-retrain",
]
_add_bullet_list(slide, Inches(1.05), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in steps_ml], font_size=13, color=LIGHT_GRAY)

_add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(7.0), Inches(3.35), Inches(0.05), Inches(0.5), ACCENT_GOLD)
_add_text_box(slide, Inches(7.25), Inches(3.3), Inches(5.1), Inches(0.5),
              "Notebooks", font_size=18, color=ACCENT_GOLD, bold=True)
notebooks_ml = [
    "01_feature_engineering.py",
    "02_model_training_hpo_optuna.py",
    "03a_create_deployment_job.py",
    "03b_from_notebook_to_models_in_uc.py",
    "04a_challenger_validation.py",
    "04b_challenger_approval.py",
    "05_batch_inference.py",
    "06_serve_features_and_model.py",
    "07_model_monitoring.py",
    "08_drift_detection.py",
]
_add_bullet_list(slide, Inches(7.25), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in notebooks_ml], font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 – Deep Dive: AI Engineering
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "AI Engineering",
                "Agent Bricks  |  Genie Spaces  |  UC Functions  |  Model Serving  |  75 min")

_add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(0.8), Inches(3.35), Inches(0.05), Inches(0.5), GREEN)
_add_text_box(slide, Inches(1.05), Inches(3.3), Inches(5.1), Inches(0.5),
              "Key Topics", font_size=18, color=GREEN, bold=True)
steps_ai = [
    "Agent Bricks: Knowledge Assistants & Multi-Agent Supervisors",
    "Genie Spaces: self-service SQL analytics via natural language",
    "UC Functions: create tool functions for agent orchestration",
    "Model Serving: deploy agents & models to endpoints",
    "ResponsesAgent + LangGraph custom agent pattern",
    "End-to-end: approve, analyze, and shop lender rates",
]
_add_bullet_list(slide, Inches(1.05), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in steps_ai], font_size=13, color=LIGHT_GRAY)

_add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(7.0), Inches(3.35), Inches(0.05), Inches(0.5), GREEN)
_add_text_box(slide, Inches(7.25), Inches(3.3), Inches(5.1), Inches(0.5),
              "Notebooks & Components", font_size=18, color=GREEN, bold=True)
notebooks_ai = [
    "09a_generate_lender_programs.py",
    "09b_setup_agent_bricks.py",
    "src/agent/agent.py         (ResponsesAgent)",
    "src/agent/test_agent.py    (test scenarios)",
    "src/agent/log_model.py     (MLflow registration)",
    "",
    "Databricks Components:",
    "  Genie Space: CDK Lending Analytics",
    "  MAS: CDK Lending Supervisor",
    "  Model Serving Endpoint (custom agent)",
]
_add_bullet_list(slide, Inches(7.25), Inches(3.9), Inches(5.1), Inches(2.8),
                 [f"  {s}" for s in notebooks_ai], font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 – Architecture Summary
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Architecture Summary")

layers = [
    ("Data Ingestion", "Auto Loader, ai_parse_document, ai_query"),
    ("Data Pipeline", "Spark Declarative Pipeline (Bronze/Silver/Gold)"),
    ("Feature Store", "Feature table + on-demand UC functions"),
    ("ML Training", "Optuna HPO + pyfunc wrapper with business rules"),
    ("Serving", "Model Serving endpoint + Online Feature Store"),
    ("Monitoring", "Lakehouse Monitoring + auto-retrain on drift"),
    ("Analytics", "Genie Space for self-service SQL"),
    ("Rate Shopping", "ResponsesAgent + LangGraph + UC Functions"),
    ("Orchestration", "Multi-Agent Supervisor (MAS)"),
]

y = Inches(3.2)
for i, (layer, tech) in enumerate(layers):
    color = [ACCENT_RED, ACCENT_GOLD, BLUE, GREEN, WHITE, MED_GRAY, BLUE, GREEN, ACCENT_RED][i]
    _add_card(slide, Inches(0.8), y, Inches(11.7), Inches(0.42))
    _accent_bar(slide, Inches(1.0), y + Inches(0.06), Inches(0.05), Inches(0.3), color)
    _add_text_box(slide, Inches(1.3), y + Inches(0.02), Inches(3.5), Inches(0.38),
                  layer, font_size=14, color=WHITE, bold=True)
    _add_text_box(slide, Inches(5.0), y + Inches(0.02), Inches(7.2), Inches(0.38),
                  tech, font_size=13, color=LIGHT_GRAY)
    y += Inches(0.46)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 – Conclusion & Resources
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(slide)

# Title
_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(0.08), Inches(0.8), ACCENT_RED)
_add_text_box(slide, Inches(0.9), Inches(0.9), Inches(11), Inches(1),
              "Conclusion & Next Steps", font_size=36, color=WHITE, bold=True)
_add_text_box(slide, Inches(0.9), Inches(1.8), Inches(11), Inches(0.5),
              "Continue your Databricks learning journey with these resources",
              font_size=18, color=LIGHT_GRAY)

# ── Left column: Databricks Training & Certification ──
_add_card(slide, Inches(0.6), Inches(2.6), Inches(5.8), Inches(4.4))
_accent_bar(slide, Inches(0.8), Inches(2.75), Inches(0.05), Inches(0.5), ACCENT_RED)
_add_text_box(slide, Inches(1.05), Inches(2.7), Inches(5.1), Inches(0.5),
              "Training & Certification", font_size=18, color=ACCENT_RED, bold=True)
training_items = [
    "Databricks Academy  (academy.databricks.com)",
    "  Free self-paced courses for all skill levels",
    "",
    "Databricks Certified Data Engineer Associate",
    "  Delta Lake, ELT, pipelines, governance",
    "",
    "Databricks Certified ML Professional",
    "  MLflow, Feature Store, model deployment",
    "",
    "Databricks Certified GenAI Engineer Associate",
    "  RAG, agents, compound AI systems",
    "",
    "Databricks Partner Connect  (partner workshops)",
    "  Hands-on labs with Databricks SAs",
]
_add_bullet_list(slide, Inches(1.05), Inches(3.25), Inches(5.1), Inches(3.5),
                 training_items, font_size=12, color=LIGHT_GRAY)

# ── Right column: Documentation & Community ──
_add_card(slide, Inches(6.8), Inches(2.6), Inches(5.8), Inches(4.4))
_accent_bar(slide, Inches(7.0), Inches(2.75), Inches(0.05), Inches(0.5), BLUE)
_add_text_box(slide, Inches(7.25), Inches(2.7), Inches(5.1), Inches(0.5),
              "Documentation & Community", font_size=18, color=BLUE, bold=True)
resource_items = [
    "docs.databricks.com",
    "  Official documentation (AWS, Azure, GCP)",
    "",
    "Databricks Community Forum",
    "  community.databricks.com",
    "",
    "GitHub: databricks/databricks-sdk-py",
    "  Python SDK, CLI, notebooks, examples",
    "",
    "Solution Accelerators",
    "  databricks.com/solutions/accelerators",
    "",
    "DATA+AI Summit (annual conference)",
    "  Sessions, workshops, & networking",
]
_add_bullet_list(slide, Inches(7.25), Inches(3.25), Inches(5.1), Inches(3.5),
                 resource_items, font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
output_path = "/Users/zach.jacobson/Desktop/github_repos/customer_repos/cdk/cdk-global/CDK_Global_Databricks_Solution.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
