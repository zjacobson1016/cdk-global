#!/usr/bin/env python3
"""Generate CDK Global Parts Invoice Processing pitch deck for Databricks."""

import ctypes.util
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Cairo libcairo fix for macOS Homebrew (before importing cairosvg)
_orig_find_library = ctypes.util.find_library
def _find_library_with_homebrew(name):
    result = _orig_find_library(name)
    if result is None and name in ("cairo", "cairo-2", "libcairo-2"):
        for path in [
            "/opt/homebrew/lib/libcairo.2.dylib",
            "/usr/local/lib/libcairo.2.dylib",
        ]:
            if os.path.exists(path):
                return path
    return result
ctypes.util.find_library = _find_library_with_homebrew

try:
    import cairosvg
    HAS_CAIROSVG = True
except (ImportError, OSError):
    HAS_CAIROSVG = False

# ── Databricks brand palette ──
BG_DARK     = RGBColor(0x1B, 0x1F, 0x23)
BG_CARD     = RGBColor(0x25, 0x2A, 0x30)
ACCENT_RED  = RGBColor(0xFF, 0x38, 0x21)
ACCENT_GOLD = RGBColor(0xFF, 0xB7, 0x2B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xBB, 0xBB, 0xBB)
MED_GRAY    = RGBColor(0x88, 0x88, 0x88)
GREEN       = RGBColor(0x00, 0xC8, 0x53)
BLUE        = RGBColor(0x42, 0xA5, 0xF5)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

ASSETS_DIR = "/Users/zach.jacobson/Desktop/github_repos/customer_repos/presentation_automation/Databricks-selected-assets"

LOGO_MAP = {
    "databricks_main": os.path.join(ASSETS_DIR, "Databricks One Lockup Full Color", "databricks-one-lockup-full-color-white.svg"),
    "unity_catalog": os.path.join(ASSETS_DIR, "Unity Catalog Lockup Full Color", "unity-catalog-lockup-full-color-white.svg"),
    "delta_lake": os.path.join(ASSETS_DIR, "logo-color-delta-lake.svg"),
    "lakeflow_pipelines": os.path.join(ASSETS_DIR, "Lakeflow Declarative Pipelines Lockup Full Color", "lakeflow-declarative-pipelines-lockup-full-color-white.svg"),
    "spark": os.path.join(ASSETS_DIR, "Apache Spark Logo", "apache-spark-logo-white-rgb.svg"),
    "mlflow": os.path.join(ASSETS_DIR, "MLflow Logo", "mlflow-logo-white-rgb.svg"),
}

_tmp_dir = tempfile.mkdtemp(prefix="db_logos_")


def _svg_to_png(svg_path, width_px=600):
    if not HAS_CAIROSVG or not os.path.exists(svg_path):
        return None
    png_path = os.path.join(_tmp_dir, os.path.basename(svg_path).replace(".svg", ".png"))
    if os.path.exists(png_path):
        return png_path
    try:
        cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=width_px)
        return png_path
    except Exception:
        return None


def _dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def _add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _horizontal_rule(slide, left, top, width, color=ACCENT_RED, height=Inches(0.04)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _section_header(slide, text, subtitle=None):
    _dark_bg(slide)
    _add_watermark(slide)
    _accent_bar(slide, Inches(0.6), Inches(1.6), Inches(0.08), Inches(0.8), ACCENT_RED)
    _add_text_box(slide, Inches(0.9), Inches(1.5), Inches(11), Inches(1),
                  text, font_size=36, color=WHITE, bold=True)
    if subtitle:
        _add_text_box(slide, Inches(0.9), Inches(2.4), Inches(11), Inches(0.7),
                      subtitle, font_size=18, color=LIGHT_GRAY)


def _add_watermark(slide):
    _add_logo(slide, "databricks_main", Inches(10.8), Inches(6.8), height=Inches(0.4))


def _add_logo(slide, logo_key, left, top, width=None, height=None):
    svg_path = LOGO_MAP.get(logo_key)
    if not svg_path:
        return None
    png_path = _svg_to_png(svg_path)
    if not png_path:
        return None
    kwargs = {}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    if not kwargs:
        kwargs["height"] = Inches(0.6)
    return slide.shapes.add_picture(png_path, left, top, **kwargs)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(slide)
_add_logo(slide, "databricks_main", Inches(0.6), Inches(0.5), height=Inches(0.7))
_horizontal_rule(slide, Inches(0.6), Inches(2.2), Inches(2.0), ACCENT_RED)

_add_text_box(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(1.2),
              "Intelligent Parts Invoice Processing — CDK Global",
              font_size=48, color=WHITE, bold=True)
_add_text_box(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(1),
              "End-to-End Automation on Databricks",
              font_size=28, color=ACCENT_GOLD, bold=True)
_add_text_box(slide, Inches(0.6), Inches(4.7), Inches(10), Inches(1),
              "Document AI  |  Medallion Architecture  |  Multi-Agent Supervisor  |  Slack Notifications",
              font_size=18, color=LIGHT_GRAY)
_add_text_box(slide, Inches(0.6), Inches(6.5), Inches(10), Inches(0.5),
              "Confidential  |  Prepared for CDK Global",
              font_size=14, color=MED_GRAY)

# Product logos along bottom
_add_logo(slide, "spark", Inches(0.6), Inches(6.7), height=Inches(0.4))
_add_logo(slide, "lakeflow_pipelines", Inches(2.2), Inches(6.7), height=Inches(0.4))
_add_logo(slide, "unity_catalog", Inches(5.0), Inches(6.7), height=Inches(0.4))
_add_logo(slide, "mlflow", Inches(7.5), Inches(6.7), height=Inches(0.4))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 – The Challenge
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "The Challenge",
                "Automotive dealerships face significant pain in parts invoice processing")

challenge_items = [
    "Automotive dealerships process hundreds of parts invoices monthly",
    "Manual 3-way matching (invoice vs PO vs receiving report) is slow and error-prone",
    "Approval routing is inconsistent and bottlenecked",
    "No visibility into supplier performance or discrepancy patterns",
    "Disconnected systems: paper PDFs, emails, spreadsheets",
]
_add_bullet_list(slide, Inches(0.9), Inches(3.2), Inches(11.5), Inches(3.5),
                 challenge_items, font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 – The Solution Overview
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "The Solution Overview",
                "End-to-end automated invoice processing on the Databricks Lakehouse")

solution_items = [
    "End-to-end automated invoice processing pipeline on Databricks",
    "AI-powered document parsing and entity extraction from PDF invoices",
    "Automated 3-way matching with intelligent approval routing",
    "Human-in-the-loop agentic approval workflow with Slack notifications",
    "Multi-Agent Supervisor for natural language interaction",
    "Real-time analytics and dashboards",
]
_add_bullet_list(slide, Inches(0.9), Inches(3.2), Inches(11.5), Inches(3.5),
                 solution_items, font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 – Architecture (Data Engineering)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Architecture: Data Engineering",
                "Medallion pipeline with AI enrichment")

_add_logo(slide, "lakeflow_pipelines", Inches(10.5), Inches(1.5), height=Inches(0.5))

# Bronze row
_add_card(slide, Inches(0.6), Inches(3.2), Inches(11.7), Inches(1.1))
_accent_bar(slide, Inches(0.8), Inches(3.35), Inches(0.05), Inches(0.8), ACCENT_RED)
_add_text_box(slide, Inches(1.0), Inches(3.3), Inches(2.5), Inches(0.4),
              "Bronze", font_size=18, color=ACCENT_RED, bold=True)
_add_text_box(slide, Inches(3.5), Inches(3.35), Inches(8.5), Inches(0.8),
              "Streaming ingestion via Spark Declarative Pipelines from UC Volumes "
              "(suppliers, POs, invoices, receiving reports, emails, PDF documents)",
              font_size=14, color=LIGHT_GRAY)

# Silver row
_add_card(slide, Inches(0.6), Inches(4.5), Inches(11.7), Inches(1.1))
_accent_bar(slide, Inches(0.8), Inches(4.65), Inches(0.05), Inches(0.8), ACCENT_GOLD)
_add_text_box(slide, Inches(1.0), Inches(4.6), Inches(2.5), Inches(0.4),
              "Silver", font_size=18, color=ACCENT_GOLD, bold=True)
_add_text_box(slide, Inches(3.5), Inches(4.65), Inches(8.5), Inches(0.8),
              "Data cleaning + AI enrichment: ai_parse_document() for PDF text extraction, "
              "ai_query() with Llama 3.3 70B for structured entity extraction",
              font_size=14, color=LIGHT_GRAY)

# Gold row
_add_card(slide, Inches(0.6), Inches(5.8), Inches(11.7), Inches(1.1))
_accent_bar(slide, Inches(0.8), Inches(5.95), Inches(0.05), Inches(0.8), GREEN)
_add_text_box(slide, Inches(1.0), Inches(5.9), Inches(2.5), Inches(0.4),
              "Gold", font_size=18, color=GREEN, bold=True)
_add_text_box(slide, Inches(3.5), Inches(5.95), Inches(8.5), Inches(0.8),
              "Materialized view: 3-way match (Invoice PDF ↔ PO ↔ Receiving Report), variance calculation, "
              "classification (STANDARD, DISCREPANCY, UNMATCHED, RECEIVING_ISSUE), and approval routing",
              font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 – Architecture (AI/Agent Layer)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Architecture: AI & Agent Layer",
                "LangGraph agent with UC Functions and Slack routing")

# Left column
_add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(0.8), Inches(3.35), Inches(0.05), Inches(0.5), ACCENT_RED)
_add_text_box(slide, Inches(1.05), Inches(3.3), Inches(5.2), Inches(0.5),
              "Invoice Processing Agent", font_size=18, color=ACCENT_RED, bold=True)
agent_items = [
    "LangGraph Agent as MLflow ResponsesAgent on Model Serving",
    "Intent Router → 7 workflows: process_invoice, approve_invoice, reject_invoice, "
    "escalate_invoice, check_status, my_approvals, general_query",
    "8 UC Functions (read-only) + 4 Python tools (writes via Lakebase)",
    "Slack notifications by route: #service-approvals, #parts-approvals, "
    "#gm-approvals, #invoice-exceptions",
]
_add_bullet_list(slide, Inches(1.05), Inches(3.9), Inches(5.2), Inches(2.8),
                 agent_items, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 – Multi-Agent Supervisor
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Multi-Agent Supervisor",
                "Orchestrating specialized agents for unified invoice operations")

# MAS card (top)
_add_card(slide, Inches(2.5), Inches(3.0), Inches(8.3), Inches(1.0), ACCENT_RED)
_add_text_box(slide, Inches(2.8), Inches(3.1), Inches(7.7), Inches(0.8),
              "Databricks Multi-Agent Supervisor",
              font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(2.8), Inches(3.6), Inches(7.7), Inches(0.5),
              "Unified conversational interface for invoice operations",
              font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# Two agent cards
card_y = Inches(4.3)
card_w = Inches(5.8)
gap = Inches(0.4)

_add_card(slide, Inches(0.6), card_y, card_w, Inches(2.5))
_accent_bar(slide, Inches(0.8), card_y + Inches(0.15), Inches(0.06), Inches(0.5), ACCENT_RED)
_add_text_box(slide, Inches(1.0), card_y + Inches(0.1), card_w - Inches(0.5), Inches(0.5),
              "Invoice Processing Agent", font_size=16, color=ACCENT_GOLD, bold=True)
_add_text_box(slide, Inches(1.0), card_y + Inches(0.65), card_w - Inches(0.5), Inches(1.7),
              "Structured operations: processing, approvals, rejections, escalations, "
              "status checks, supplier analytics",
              font_size=13, color=LIGHT_GRAY)

_add_card(slide, Inches(6.8), card_y, card_w, Inches(2.5))
_accent_bar(slide, Inches(7.0), card_y + Inches(0.15), Inches(0.06), Inches(0.5), BLUE)
_add_text_box(slide, Inches(7.2), card_y + Inches(0.1), card_w - Inches(0.5), Inches(0.5),
              "Genie Space (Invoice Approval Analytics)", font_size=16, color=ACCENT_GOLD, bold=True)
genie_items = [
    "Ad-hoc SQL via natural language",
    '"How many invoices pending for the parts director?"',
    '"Show vendor discrepancy rates"',
    "AI/BI Dashboard + metric views",
    "Flask conversation UI",
]
_add_bullet_list(slide, Inches(7.2), card_y + Inches(0.65), card_w - Inches(0.5), Inches(1.7),
                 genie_items, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 – Key Capabilities & Business Value
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Key Capabilities & Business Value",
                "Measurable impact on efficiency, visibility, and compliance")

capabilities = [
    "Automated 3-way matching eliminates manual reconciliation",
    "AI document parsing handles unstructured PDF invoices",
    "Intelligent routing reduces approval bottlenecks (AUTO_APPROVED for clean matches, escalation paths for discrepancies)",
    "Real-time Slack notifications keep approvers informed",
    "Natural language queries for instant invoice insights",
    "Supplier performance analytics for vendor management",
    "Full audit trail in Lakebase for compliance",
]
_add_bullet_list(slide, Inches(0.9), Inches(3.2), Inches(11.5), Inches(3.5),
                 capabilities, font_size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 – Technology Stack
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Technology Stack",
                "End-to-end Databricks platform capabilities")

stack_rows = [
    ("Data Generation", "Faker, FPDF2, PySpark"),
    ("Ingestion", "Spark Declarative Pipelines, Auto Loader"),
    ("AI Enrichment", "ai_parse_document, ai_query (Llama 3.3 70B)"),
    ("Matching & Routing", "SQL materialized views, rule-based logic"),
    ("Agent Framework", "LangGraph, MLflow ResponsesAgent, UCFunctionToolkit"),
    ("Approval State", "Lakebase Provisioned (PostgreSQL)"),
    ("Notifications", "Slack Bolt SDK, Block Kit"),
    ("Orchestration", "Databricks Multi-Agent Supervisor"),
    ("Analytics", "Genie Space, AI/BI Metric Views"),
    ("Deployment", "Databricks Asset Bundles, Model Serving"),
]

y = Inches(3.2)
for layer, tech in stack_rows:
    _add_card(slide, Inches(0.8), y, Inches(11.7), Inches(0.42))
    _accent_bar(slide, Inches(1.0), y + Inches(0.06), Inches(0.05), Inches(0.3), ACCENT_RED)
    _add_text_box(slide, Inches(1.3), y + Inches(0.02), Inches(3.5), Inches(0.38),
                  layer, font_size=14, color=WHITE, bold=True)
    _add_text_box(slide, Inches(5.0), y + Inches(0.02), Inches(7.2), Inches(0.38),
                  tech, font_size=13, color=LIGHT_GRAY)
    y += Inches(0.46)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 – Deployment & Next Steps
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_section_header(slide, "Deployment & Next Steps",
                "Databricks Asset Bundles enable repeatable, multi-environment deployments")

# Left - Deployment
_add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(0.8), Inches(3.35), Inches(0.05), Inches(0.5), ACCENT_RED)
_add_text_box(slide, Inches(1.05), Inches(3.3), Inches(5.2), Inches(0.5),
              "Deployment", font_size=18, color=ACCENT_RED, bold=True)
deploy_items = [
    "Databricks Asset Bundles (DABs) with dev/prod targets",
    "End-to-end orchestration job:",
    "  generate_data → generate_pdfs → run_pipeline → create_uc_functions → log_model → deploy_agent",
]
_add_bullet_list(slide, Inches(1.05), Inches(3.9), Inches(5.2), Inches(2.8),
                 deploy_items, font_size=14, color=LIGHT_GRAY)

# Right - Next Steps
_add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8))
_accent_bar(slide, Inches(7.0), Inches(3.35), Inches(0.05), Inches(0.5), GREEN)
_add_text_box(slide, Inches(7.25), Inches(3.3), Inches(5.1), Inches(0.5),
              "Next Steps", font_size=18, color=GREEN, bold=True)
next_items = [
    "Production rollout",
    "Integration with CDK DMS",
    "Expanded supplier coverage",
    "Additional approval workflows",
]
_add_bullet_list(slide, Inches(7.25), Inches(3.9), Inches(5.1), Inches(2.8),
                 next_items, font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 – Closing / Thank You
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(slide)

_add_logo(slide, "databricks_main", Inches(4.8), Inches(1.0), height=Inches(1.0))
_horizontal_rule(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT_RED)
_add_text_box(slide, Inches(0.6), Inches(2.9), Inches(12), Inches(1.2),
              "Transform Parts Invoice Processing at CDK Global",
              font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(0.8),
              "Ready to automate, accelerate, and gain visibility across your invoice workflow?",
              font_size=24, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.6),
              "Contact your Databricks team to get started",
              font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
_add_watermark(slide)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
output_path = "/Users/zach.jacobson/Desktop/github_repos/customer_repos/cdk/cdk-global/parts_invoice_processing/CDK_Global_Parts_Invoice_Processing.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
